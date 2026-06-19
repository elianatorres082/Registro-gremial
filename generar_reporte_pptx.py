"""
Generador de reporte PPT en Python puro (python-pptx).
Sin dependencia de Node.js — funciona en Streamlit Cloud.

Uso desde app.py:
    from generar_reporte_pptx import generar_reporte_bytes
    ppt_bytes = generar_reporte_bytes(df, periodo_label)
"""

import io
import re
from datetime import date
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colores ───────────────────────────────────────────────────────────────────
def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY   = "1A2C5B"
BLUE   = "2E4A9E"
RED    = "C0392B"
ORANGE = "E67E22"
GREEN  = "27AE60"
PURPLE = "8E44AD"
GRAY   = "64748B"
LGRAY  = "F8FAFC"
BGRAY  = "E2E8F0"
WHITE  = "FFFFFF"
BLACK  = "1E293B"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=11, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox

def add_slide_header(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Header bar
    add_rect(slide, 0, 0, 10, 0.75, NAVY)
    add_text(slide, title.upper(), 0.3, 0.05, 9.4, 0.65,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.75, 9.4, 0.28,
                 size=8.5, color=GRAY)
    return slide

def add_kpi_card(slide, x, y, w, h, label, value, sub="", border_color=BLUE):
    add_rect(slide, x, y, w, h, WHITE, BGRAY, Pt(0.75))
    add_rect(slide, x, y, w, 0.05, border_color)
    add_text(slide, label.upper(), x+0.1, y+0.08, w-0.2, 0.2,
             size=7, color=GRAY)
    add_text(slide, str(value), x+0.1, y+0.28, w-0.2, 0.5,
             size=28, bold=True, color=BLACK)
    if sub:
        add_text(slide, sub, x+0.1, y+0.78, w-0.2, 0.2, size=7.5, color=GRAY)

def pct(n, total):
    return round(n/total*100, 1) if total > 0 else 0

def fmt_pct(v):
    return f"{v:.1f}%"

MESES_ES = {1:"Enero",2:"Febrero",3:"Marzo",4:"Abril",5:"Mayo",6:"Junio",
            7:"Julio",8:"Agosto",9:"Septiembre",10:"Octubre",11:"Noviembre",12:"Diciembre"}

# ── Clasificar desvío ──────────────────────────────────────────────────────────
DESVIOS_NORM = [
    "no cumple","excede -5min","excede 5/10 min","excede + 10/20min",
    "supera + 1/2hs","supera semana completa","supera semanal",
    "genera pdl","cambia de turno","no realiza actividad gremial",
]

def es_desvio(val):
    return str(val).strip().lower() in DESVIOS_NORM

def clasificar_motivo(row):
    vals = " ".join([
        str(row.get("Movilidad Gremial x Semana x Dia","")),
        str(row.get("Motivo Excedencia","")),
        str(row.get("Licencia","")),
        str(row.get("LLT/Ausencia Inj","")),
        str(row.get("Observaciones Extras","")),
        str(row.get("Motivo (Detallar desvios)","")),
    ]).lower()
    if any(x in vals for x in ["llt","charla 5","charla5"]):
        return "LLT / Charla"
    if any(x in vals for x in ["art","médica","medica","enfermedad"]):
        return "Ausencia / Licencia Medica o ART"
    if any(x in vals for x in ["incumplimiento","injustificad","inj","no cumple"]):
        return "Incumplimiento / Ausencia sin justificar"
    if any(x in vals for x in ["excede","supera","exceso"]):
        return "Exceso en Movilidad Gremial"
    if any(x in vals for x in ["moviliz","asamblea","reunion","gremial","actividad"]):
        return "Movilizacion / Actividad Gremial"
    if any(x in vals for x in ["cambia","pdl","no realiza"]):
        return "Actividad Gremial (otras)"
    return "Otros"

def parsear_fecha(val):
    s = str(val).strip()
    if s in ("","nan","NaT","None"): return None
    m = re.search(r"(\d{1,4})[/\-\.](\d{1,2})[/\-\.](\d{2,4})", s)
    if not m: return None
    a,b,c = m.group(1), m.group(2), m.group(3)
    try:
        if len(a)==4: return date(int(a),int(b),int(c))
        yr = int(c) if len(c)==4 else 2000+int(c)
        return date(yr,int(b),int(a))
    except: return None

# ── Calcular datos ─────────────────────────────────────────────────────────────
def calcular(df):
    rows = df.to_dict("records")
    for r in rows:
        r["_fecha"] = parsear_fecha(r.get("Fecha",""))
        r["_desvio"] = es_desvio(r.get("Movilidad Gremial x Semana x Dia",""))
        r["_motivo"] = clasificar_motivo(r) if r["_desvio"] else ""

    total = len(rows)
    desvios = [r for r in rows if r["_desvio"]]
    n_dev = len(desvios)
    tasa = pct(n_dev, total)

    # Tendencia mensual
    por_mes = {}
    for r in rows:
        if not r["_fecha"]: continue
        key = (r["_fecha"].year, r["_fecha"].month)
        if key not in por_mes: por_mes[key] = {"total":0,"dev":0}
        por_mes[key]["total"] += 1
        if r["_desvio"]: por_mes[key]["dev"] += 1
    tendencia = sorted([
        {"anio":k[0],"mes":k[1],"label":MESES_ES[k[1]],"total":v["total"],
         "dev":v["dev"],"pct":pct(v["dev"],v["total"])}
        for k,v in por_mes.items()
    ], key=lambda x:(x["anio"],x["mes"]))

    # Por sector
    por_sec = {}
    for r in rows:
        s = r.get("Sector","Sin sector")
        if s not in por_sec: por_sec[s] = {"total":0,"dev":0,"motivos":{}}
        por_sec[s]["total"] += 1
        if r["_desvio"]:
            por_sec[s]["dev"] += 1
            m = r["_motivo"]
            por_sec[s]["motivos"][m] = por_sec[s]["motivos"].get(m,0)+1
    sectores = sorted([
        {"sector":k,"total":v["total"],"dev":v["dev"],
         "pct":pct(v["dev"],v["total"]),"motivos":v["motivos"]}
        for k,v in por_sec.items()
    ], key=lambda x:-x["pct"])

    # Por turno
    por_turno = {}
    for r in rows:
        t = r.get("Turno","Sin turno")
        if t not in por_turno: por_turno[t] = {"total":0,"dev":0,"motivos":{}}
        por_turno[t]["total"] += 1
        if r["_desvio"]:
            por_turno[t]["dev"] += 1
            m = r["_motivo"]
            por_turno[t]["motivos"][m] = por_turno[t]["motivos"].get(m,0)+1
    turnos = sorted([
        {"turno":k,"total":v["total"],"dev":v["dev"],
         "pct":pct(v["dev"],v["total"]),"motivos":v["motivos"]}
        for k,v in por_turno.items()
    ], key=lambda x:-x["pct"])

    # Top 5 delegados
    por_del = {}
    for r in rows:
        n = r.get("Nombre y Apellido","")
        if not n: continue
        if n not in por_del:
            por_del[n] = {"nombre":n,"sector":r.get("Sector",""),"turno":r.get("Turno",""),
                          "total":0,"dev":0,"motivos":{}}
        por_del[n]["total"] += 1
        if r["_desvio"]:
            por_del[n]["dev"] += 1
            m = r["_motivo"]
            por_del[n]["motivos"][m] = por_del[n]["motivos"].get(m,0)+1
    top5 = sorted([
        {**v,"pct":pct(v["dev"],v["total"])}
        for v in por_del.values()
    ], key=lambda x:-x["dev"])[:5]

    # Motivos globales
    motivos_g = {}
    for r in desvios:
        m = r["_motivo"]
        motivos_g[m] = motivos_g.get(m,0)+1
    motivos_global = sorted([
        {"motivo":k,"cant":v,"pct":pct(v,n_dev)}
        for k,v in motivos_g.items()
    ], key=lambda x:-x["cant"])

    # LLT
    llt_rows = [r for r in desvios if r["_motivo"]=="LLT / Charla"]
    por_del_llt = {}
    for r in llt_rows:
        n = r.get("Nombre y Apellido","")
        if n not in por_del_llt:
            por_del_llt[n] = {"nombre":n,"sector":r.get("Sector",""),"turno":r.get("Turno",""),"cant":0}
        por_del_llt[n]["cant"] += 1
    llt_top = sorted(por_del_llt.values(), key=lambda x:-x["cant"])

    return {
        "total":total,"n_dev":n_dev,"tasa":tasa,
        "tendencia":tendencia,"sectores":sectores,"turnos":turnos,
        "top5":top5,"motivos_global":motivos_global,
        "llt_rows":llt_rows,"llt_top":llt_top
    }

# ── Generar PPT ────────────────────────────────────────────────────────────────
def generar_reporte_bytes(df, periodo_label="", n_delegados=None):
    D = calcular(df)
    if not periodo_label:
        periodo_label = f"{D['total']} registros"
    if n_delegados is None:
        n_delegados = df["Legajo"].nunique() if "Legajo" in df.columns else "?"

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════
    # SLIDE 1 — PORTADA
    # ══════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 10, 5.625, NAVY)
    add_rect(s1, 0, 4.1, 10, 1.525, BLUE)
    add_text(s1, "GESTIÓN DE PRESENCIA Y MOVILIDAD GREMIAL",
             0.5, 0.9, 9, 0.7, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s1, "Análisis — Evolución del Proceso",
             0.5, 1.7, 9, 0.45, size=14, italic=True, color="A8BFFF", align=PP_ALIGN.CENTER)
    add_rect(s1, 3.5, 2.3, 3, 0.04, ORANGE)
    add_text(s1, f"Período: {periodo_label}  |  {n_delegados} Delegados  |  2 hs Movilidad Gremial",
             0.5, 2.45, 9, 0.4, size=10, color="A8BFFF", align=PP_ALIGN.CENTER)
    add_text(s1, "Toyota Argentina S.A.",
             0.5, 4.25, 9, 0.35, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s1, "Recursos Humanos — Relaciones Laborales",
             0.5, 4.6, 9, 0.3, size=9, color="A8BFFF", align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 2 — INDICADORES CLAVE
    # ══════════════════════════════════════════════════
    s2 = add_slide_header(prs, "Contexto — Indicadores Claves", periodo_label)
    n_cumple = D["total"] - D["n_dev"]
    tasa_color = RED if D["tasa"]>15 else ORANGE if D["tasa"]>10 else GREEN
    kpis = [
        ("Registros Analizados", f"{D['total']:,}".replace(",","."), periodo_label, BLUE),
        ("Total de Desvios", str(D["n_dev"]), "", RED),
        ("Tasa de Desvio Global", fmt_pct(D["tasa"]), f"sobre {D['total']} registros", tasa_color),
    ]
    kw = 2.1
    for i,(lbl,val,sub,col) in enumerate(kpis):
        add_kpi_card(s2, 0.35+i*(kw+0.18), 1.05, kw, 1.25, lbl, val, sub, col)

    # Card destacada de Registros Cumple (igual que el reporte manual)
    add_rect(s2, 2.6, 2.45, 4.5, 0.95, GREEN)
    add_text(s2, f"{n_cumple:,}".replace(",","."), 2.6, 2.52, 4.5, 0.5,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, "Registros Cumple", 2.6, 3.0, 4.5, 0.3,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Tendencia en texto (debajo de la card de Cumple)
    add_text(s2, "EVOLUCION MENSUAL", 0.35, 3.5, 9, 0.25, size=8, bold=True, color=NAVY)
    if D["tendencia"]:
        col_w_t = 9.0 / len(D["tendencia"])
        bar_base_y = 3.8
        bar_max_h = 0.9
        for i, t in enumerate(D["tendencia"]):
            tc = RED if t["pct"]>=20 else ORANGE if t["pct"]>=12 else GREEN
            bar_h = min(t["pct"]/30 * bar_max_h, bar_max_h)
            add_rect(s2, 0.35+i*col_w_t, bar_base_y+bar_max_h-bar_h, col_w_t-0.08, bar_h, tc)
            add_text(s2, t["label"][:3], 0.35+i*col_w_t, bar_base_y+bar_max_h+0.05, col_w_t-0.08, 0.22,
                     size=7, color=GRAY, align=PP_ALIGN.CENTER)
            add_text(s2, fmt_pct(t["pct"]), 0.35+i*col_w_t, bar_base_y+bar_max_h+0.27, col_w_t-0.08, 0.22,
                     size=7.5, bold=True, color=tc, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 3 — TENDENCIA MENSUAL
    # ══════════════════════════════════════════════════
    s3 = add_slide_header(prs, "Tendencia Mensual de Desvios", periodo_label)
    if D["tendencia"]:
        max_pct_t = max(t["pct"] for t in D["tendencia"]) or 1
        bar_area_h = 3.5
        bar_area_y = 1.1
        bar_w = min(7.0/len(D["tendencia"])-0.15, 0.9)
        x_start = 0.5
        gap = (7.0 - len(D["tendencia"])*bar_w) / max(len(D["tendencia"])-1,1)

        for i,t in enumerate(D["tendencia"]):
            tc = RED if t["pct"]>=20 else ORANGE if t["pct"]>=12 else GREEN
            bar_h = (t["pct"]/max_pct_t)*bar_area_h
            bx = x_start + i*(bar_w+gap)
            by = bar_area_y + bar_area_h - bar_h
            add_rect(s3, bx, by, bar_w, bar_h, tc)
            add_text(s3, fmt_pct(t["pct"]), bx, by-0.3, bar_w, 0.25,
                     size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)
            add_text(s3, t["label"][:3].upper(), bx, bar_area_y+bar_area_h+0.05, bar_w, 0.25,
                     size=8, color=GRAY, align=PP_ALIGN.CENTER)

        # Cards de meses a la derecha
        card_h3 = (4.3-0.1*(len(D["tendencia"])-1)) / len(D["tendencia"])
        for i,t in enumerate(D["tendencia"]):
            tc = RED if t["pct"]>=20 else ORANGE if t["pct"]>=12 else GREEN
            yy = 1.1 + i*(card_h3+0.1)
            add_rect(s3, 7.8, yy, 1.9, card_h3, LGRAY, BGRAY, Pt(0.5))
            add_text(s3, t["label"], 7.85, yy+0.04, 1.8, 0.25, size=8, bold=True, color=tc)
            add_text(s3, fmt_pct(t["pct"]), 7.85, yy+0.28, 1.8, 0.28, size=14, bold=True, color=tc)
            add_text(s3, f"{t['dev']} dev.", 7.85, yy+0.56, 1.8, 0.2, size=7.5, color=GRAY)

    # ══════════════════════════════════════════════════
    # SLIDE 4 — DESVIOS POR SECTOR
    # ══════════════════════════════════════════════════
    s4 = add_slide_header(prs, "Desvios por Sector", periodo_label)
    headers_s4 = ["Sector","Registros","Desvios","% Desvio","Nivel"]
    col_xs = [0.3, 4.1, 5.2, 6.2, 7.3]
    col_ws = [3.7, 1.0, 1.0, 1.0, 2.4]
    row_h  = 0.31
    # Header
    add_rect(s4, 0.3, 1.0, 9.4, row_h, NAVY)
    for j,(hdr,cx,cw) in enumerate(zip(headers_s4,col_xs,col_ws)):
        add_text(s4, hdr, cx+0.05, 1.0, cw-0.1, row_h, size=8.5, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)
    # Rows
    for i,sec in enumerate(D["sectores"]):
        yy = 1.0 + (i+1)*row_h
        if yy > 5.0: break
        bg = "FEE2E2" if sec["pct"]>=15 else "FEF3C7" if sec["pct"]>=10 else "DCFCE7" if sec["pct"]>0 else "F8FAFC"
        add_rect(s4, 0.3, yy, 9.4, row_h, "F8FAFC" if i%2==0 else WHITE, BGRAY, Pt(0.3))
        tc2 = RED if sec["pct"]>=15 else ORANGE if sec["pct"]>=10 else GREEN
        nivel = "Alto" if sec["pct"]>=15 else "Medio" if sec["pct"]>=10 else "Bajo"
        add_text(s4, sec["sector"], col_xs[0]+0.05, yy+0.03, col_ws[0]-0.1, row_h-0.05, size=8.5, color=BLACK)
        add_text(s4, str(sec["total"]), col_xs[1]+0.05, yy+0.03, col_ws[1]-0.1, row_h-0.05, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s4, str(sec["dev"]),   col_xs[2]+0.05, yy+0.03, col_ws[2]-0.1, row_h-0.05, size=8.5, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(s4, fmt_pct(sec["pct"]), col_xs[3]+0.05, yy+0.03, col_ws[3]-0.1, row_h-0.05, size=9, bold=True, color=tc2, align=PP_ALIGN.CENTER)
        add_rect(s4, col_xs[4]+0.1, yy+0.05, col_ws[4]-0.2, row_h-0.1, bg)
        add_text(s4, nivel, col_xs[4]+0.1, yy+0.03, col_ws[4]-0.2, row_h-0.05, size=8.5, bold=True, color=tc2, align=PP_ALIGN.CENTER)

    # Total
    yy_tot = 1.0 + (len(D["sectores"])+1)*row_h
    if yy_tot < 5.1:
        add_rect(s4, 0.3, yy_tot, 9.4, row_h, "F1F5F9", BGRAY, Pt(0.5))
        add_text(s4, "TOTAL GENERAL", col_xs[0]+0.05, yy_tot+0.03, col_ws[0]-0.1, row_h-0.05, size=8.5, bold=True, color=BLACK)
        add_text(s4, str(D["total"]), col_xs[1]+0.05, yy_tot+0.03, col_ws[1]-0.1, row_h-0.05, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(s4, str(D["n_dev"]), col_xs[2]+0.05, yy_tot+0.03, col_ws[2]-0.1, row_h-0.05, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(s4, fmt_pct(D["tasa"]), col_xs[3]+0.05, yy_tot+0.03, col_ws[3]-0.1, row_h-0.05, size=8.5, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 5 — MOTIVOS POR SECTOR (top 6)
    # ══════════════════════════════════════════════════
    s5 = add_slide_header(prs, "Motivos por Sector — 6 sectores mas criticos", periodo_label)
    top6 = D["sectores"][:6]
    for i,sec in enumerate(top6):
        col = i%3; row2 = i//3
        cx = 0.2 + col*3.25
        cy = 1.05 + row2*2.2
        tc3 = RED if sec["pct"]>=15 else ORANGE if sec["pct"]>=10 else GREEN
        add_rect(s5, cx, cy, 3.1, 2.0, LGRAY, BGRAY, Pt(0.5))
        add_text(s5, f"{sec['sector'].replace(' Division','')}  {fmt_pct(sec['pct'])}",
                 cx+0.1, cy+0.08, 2.9, 0.28, size=8.5, bold=True, color=tc3)
        motivos_s = sorted(sec["motivos"].items(), key=lambda x:-x[1])[:4]
        for mi,(mot,cnt) in enumerate(motivos_s):
            pm = pct(cnt, sec["dev"])
            mc = mot.replace("Movilizacion / Actividad Gremial","Movilizacion")  \
                    .replace("Exceso en Movilidad Gremial","Exceso MG")  \
                    .replace("Ausencia / Licencia Medica o ART","Lic. Medica/ART")  \
                    .replace("Incumplimiento / Ausencia sin justificar","Incumplimiento")  \
                    .replace("Actividad Gremial (otras)","Act. Gremial otras")
            bar_w2 = 2.5
            add_rect(s5, cx+0.1, cy+0.4+mi*0.35, bar_w2, 0.09, BGRAY)
            if pm>0: add_rect(s5, cx+0.1, cy+0.4+mi*0.35, bar_w2*(pm/100), 0.09, tc3)
            add_text(s5, f"{mc}  {fmt_pct(pm)}", cx+0.1, cy+0.5+mi*0.35, 2.9, 0.22, size=7.5, color=BLACK)

    # ══════════════════════════════════════════════════
    # SLIDE 6 — DESVIOS POR TURNO
    # ══════════════════════════════════════════════════
    s6 = add_slide_header(prs, "Desvios por Turno", periodo_label)
    if D["turnos"]:
        bar_area_h6 = 3.8
        max_pct6 = max(t["pct"] for t in D["turnos"]) or 1
        bar_w6 = min(3.5/len(D["turnos"])-0.15, 0.7)
        gap6 = (3.5-len(D["turnos"])*bar_w6) / max(len(D["turnos"])-1,1)
        for i,t in enumerate(D["turnos"]):
            tc6 = RED if t["pct"]>=15 else ORANGE if t["pct"]>=10 else GREEN
            bh6 = (t["pct"]/max_pct6)*bar_area_h6
            bx6 = 0.4 + i*(bar_w6+gap6)
            by6 = 1.1 + bar_area_h6 - bh6
            add_rect(s6, bx6, by6, bar_w6, bh6, tc6)
            add_text(s6, fmt_pct(t["pct"]), bx6, by6-0.3, bar_w6, 0.25, size=9, bold=True, color=tc6, align=PP_ALIGN.CENTER)
            add_text(s6, t["turno"], bx6, 1.1+bar_area_h6+0.05, bar_w6, 0.25, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)

        # Cards motivos por turno
        card_h6 = (4.3-0.08*(len(D["turnos"])-1)) / len(D["turnos"])
        for i,t in enumerate(D["turnos"]):
            tc6 = RED if t["pct"]>=15 else ORANGE if t["pct"]>=10 else GREEN
            yy6 = 1.05 + i*(card_h6+0.08)
            add_rect(s6, 4.5, yy6, 5.2, card_h6, LGRAY, BGRAY, Pt(0.5))
            add_text(s6, f"TURNO {t['turno']}  {fmt_pct(t['pct'])}", 4.6, yy6+0.04, 5.0, 0.25, size=9, bold=True, color=tc6)
            mots6 = sorted(t["motivos"].items(), key=lambda x:-x[1])[:3]
            for mi2,(mot6,c6) in enumerate(mots6):
                pm6 = pct(c6,t["dev"])
                mc6 = mot6.replace("Movilizacion / Actividad Gremial","Movilizacion")  \
                          .replace("Exceso en Movilidad Gremial","Exceso MG")  \
                          .replace("Ausencia / Licencia Medica o ART","Lic. Med/ART")  \
                          .replace("Incumplimiento / Ausencia sin justificar","Incumplimiento")
                add_text(s6, f"{mc6}  {fmt_pct(pm6)}", 4.6, yy6+0.28+mi2*0.2, 5.0, 0.18, size=7.5, color=BLACK)

    # ══════════════════════════════════════════════════
    # SLIDE 7 — TOP 5 DELEGADOS
    # ══════════════════════════════════════════════════
    s7 = add_slide_header(prs, "Top 5 Delegados con Mayor Indice de Desvios", periodo_label)
    colors5 = [RED, PURPLE, ORANGE, BLUE, GREEN]
    cw5 = 1.72
    for i,d in enumerate(D["top5"]):
        cx5 = 0.3 + i*(cw5+0.1)
        col5 = colors5[i]
        add_rect(s7, cx5, 1.05, cw5, 4.1, WHITE, BGRAY, Pt(0.75))
        add_rect(s7, cx5+cw5/2-0.22, 1.15, 0.44, 0.44, col5)
        add_text(s7, f"#{i+1}", cx5+cw5/2-0.22, 1.15, 0.44, 0.44, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        apellido = " ".join(d["nombre"].split()[-2:]) if len(d["nombre"].split())>1 else d["nombre"]
        add_text(s7, apellido, cx5+0.05, 1.68, cw5-0.1, 0.4, size=9, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(s7, d["sector"].replace(" Division",""), cx5+0.05, 2.08, cw5-0.1, 0.25, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s7, str(d["dev"]), cx5+0.05, 2.35, cw5-0.1, 0.55, size=30, bold=True, color=col5, align=PP_ALIGN.CENTER)
        add_text(s7, "desvios", cx5+0.05, 2.88, cw5-0.1, 0.22, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s7, fmt_pct(d["pct"]), cx5+0.05, 3.1, cw5-0.1, 0.28, size=12, bold=True, color=col5, align=PP_ALIGN.CENTER)
        top_m = sorted(d["motivos"].items(), key=lambda x:-x[1])
        if top_m:
            mc7 = top_m[0][0].replace("Movilizacion / Actividad Gremial","Movilizacion")  \
                              .replace("Exceso en Movilidad Gremial","Exceso MG")  \
                              .replace("Ausencia / Licencia Medica o ART","Lic. Medica/ART")
            add_text(s7, f"Principal:\n{mc7}", cx5+0.05, 3.45, cw5-0.1, 0.55, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 8 — FOCALIZAR EL PROBLEMA
    # ══════════════════════════════════════════════════
    s8 = add_slide_header(prs, "Focalizar el Problema", periodo_label)
    mov_g = next((m for m in D["motivos_global"] if "oviliz" in m["motivo"]),{"pct":0})
    secs_altos = sum(1 for s in D["sectores"] if s["pct"]>=15)
    atipicos = [d for d in D["top5"] if (d["motivos"].get("Ausencia / Licencia Medica o ART",0)/d["dev"])>0.3] if D["top5"] and D["n_dev"]>0 else []

    bloques8 = [
        ("HOMOGENEO\n(Transversal)","Institucional",
         f"La Movilizacion/Actividad Gremial ({fmt_pct(mov_g['pct'])}) afecta a todos los sectores.\nRequiere decision politica.",
         f"{secs_altos}/{len(D['sectores'])} sectores", BLUE, "EEF2FF"),
        ("FOCALIZADO\n(Por sector)","Operativo",
         f"Exceso MG y ausentismo en sectores criticos.\nEstos SI son abordables.",
         f"{min(3,secs_altos)} sectores mas afectados", ORANGE, "FEF3C7"),
        ("FOCALIZADO\n(Por persona)","Individuales",
         f"{len(atipicos) if atipicos else len(D['top5'])} delegados con perfiles que requieren seguimiento individual.",
         f"{len(atipicos) if atipicos else len(D['top5'])} delegados criticos", RED, "FEE2E2"),
    ]
    for i,(tipo,titulo,texto,stat,col8,bg8) in enumerate(bloques8):
        cx8 = 0.3 + i*3.2
        add_rect(s8, cx8, 1.0, 3.1, 4.35, bg8, col8, Pt(1.5))
        add_text(s8, tipo, cx8+0.1, 1.08, 2.9, 0.55, size=9.5, bold=True, color=col8, align=PP_ALIGN.CENTER)
        add_text(s8, titulo, cx8+0.1, 1.65, 2.9, 0.28, size=8.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s8, texto, cx8+0.1, 2.0, 2.9, 1.5, size=8.5, color=BLACK)
        add_rect(s8, cx8+0.3, 3.85, 2.5, 0.6, col8)
        add_text(s8, stat, cx8+0.3, 3.9, 2.5, 0.55, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 9 — MOTIVOS GLOBALES vs TOP 5
    # ══════════════════════════════════════════════════
    s9 = add_slide_header(prs, "Analisis entre Motivos Globales con el Top 5", periodo_label)
    add_text(s9, f"GLOBAL ({D['n_dev']} desvios)", 0.3, 1.05, 3.8, 0.25, size=8.5, bold=True, color=NAVY)
    bar_max = 3.3
    for i,m in enumerate(D["motivos_global"][:7]):
        yy9 = 1.38 + i*0.51
        add_rect(s9, 0.3, yy9, bar_max, 0.16, BGRAY)
        if m["pct"]>0: add_rect(s9, 0.3, yy9, bar_max*(m["pct"]/100), 0.16, BLUE)
        mc9 = m["motivo"].replace("Movilizacion / Actividad Gremial","Movilizacion")  \
                         .replace("Exceso en Movilidad Gremial","Exceso MG")  \
                         .replace("Ausencia / Licencia Medica o ART","Lic. Med/ART")  \
                         .replace("Incumplimiento / Ausencia sin justificar","Incumplimiento")
        add_text(s9, f"{fmt_pct(m['pct'])}  {mc9}", 0.3, yy9+0.18, 3.8, 0.24, size=8, color=BLACK)

    icons9 = ["(!!)","(!!)","(MIX)","(OK)","(OK)"]
    for i,d in enumerate(D["top5"]):
        col9 = colors5[i]
        yy9b = 1.05 + i*0.87
        add_rect(s9, 4.3, yy9b, 5.4, 0.78, LGRAY, BGRAY, Pt(0.5))
        top_m9 = sorted(d["motivos"].items(), key=lambda x:-x[1])
        pct_top9 = fmt_pct(pct(top_m9[0][1],d["dev"])) if top_m9 else ""
        mc9b = top_m9[0][0].replace("Movilizacion / Actividad Gremial","Movilizacion")  \
                            .replace("Exceso en Movilidad Gremial","Exceso MG")  \
                            .replace("Ausencia / Licencia Medica o ART","Med/ART") if top_m9 else ""
        add_text(s9, f"{icons9[i] if i<len(icons9) else ''}  {d['nombre']}", 4.45, yy9b+0.05, 4.0, 0.25, size=9, bold=True, color=col9)
        add_text(s9, f"{pct_top9} {mc9b}", 4.45, yy9b+0.3, 5.1, 0.22, size=8.5, color=BLACK)
        add_text(s9, f"{d['dev']} desvios — {fmt_pct(d['pct'])}", 4.45, yy9b+0.53, 5.1, 0.2, size=8, color=GRAY)

    # ══════════════════════════════════════════════════
    # SLIDE 10 — LLT RESUMEN
    # ══════════════════════════════════════════════════
    s10 = add_slide_header(prs, "LLT / Charla 5' / Extensiones — Analisis global", periodo_label)
    n_llt = len(D["llt_rows"])
    n_del_llt = len(D["llt_top"])
    pct_llt_tot = pct(n_llt, D["total"])
    pct_llt_dev = pct(n_llt, D["n_dev"])
    kpis10 = [
        ("Total LLT/Extensiones", str(n_llt), "", ORANGE),
        ("% sobre total registros", fmt_pct(pct_llt_tot), "", BLUE),
        ("% sobre total desvios", fmt_pct(pct_llt_dev), "", RED),
        (f"Delegados afectados", f"{n_del_llt}/{n_delegados}", "", PURPLE),
    ]
    for i,(lbl,val,sub,col) in enumerate(kpis10):
        add_kpi_card(s10, 0.3+i*2.35, 1.05, 2.2, 1.15, lbl, val, sub, col)

    add_text(s10, f"LLT representa el {fmt_pct(pct_llt_dev)} del total de desvios", 0.3, 2.35, 9.4, 0.28, size=9.5, italic=True, color=GRAY)

    # Tabla LLT
    if D["llt_top"]:
        hdrs10 = ["Delegado","Sector","Turno","Desvios LLT","% del total"]
        col_xs10 = [0.3,3.2,5.4,6.6,7.8]
        col_ws10 = [2.8,2.1,1.1,1.1,1.8]
        add_rect(s10, 0.3, 2.72, 9.4, 0.28, NAVY)
        for j,(hdr,cx10,cw10) in enumerate(zip(hdrs10,col_xs10,col_ws10)):
            add_text(s10, hdr, cx10+0.04, 2.72, cw10-0.08, 0.28, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)
        for i,d in enumerate(D["llt_top"][:8]):
            yy10 = 3.0 + i*0.28
            if yy10>5.2: break
            bg10 = "F8FAFC" if i%2==0 else WHITE
            add_rect(s10, 0.3, yy10, 9.4, 0.27, bg10, BGRAY, Pt(0.3))
            add_text(s10, d["nombre"], col_xs10[0]+0.04, yy10+0.03, col_ws10[0]-0.08, 0.22, size=8.5)
            add_text(s10, d.get("sector",""), col_xs10[1]+0.04, yy10+0.03, col_ws10[1]-0.08, 0.22, size=8)
            add_text(s10, d.get("turno",""), col_xs10[2]+0.04, yy10+0.03, col_ws10[2]-0.08, 0.22, size=8, align=PP_ALIGN.CENTER)
            add_text(s10, str(d["cant"]), col_xs10[3]+0.04, yy10+0.03, col_ws10[3]-0.08, 0.22, size=8.5, bold=True, align=PP_ALIGN.CENTER)
            add_text(s10, fmt_pct(pct(d["cant"],n_llt)), col_xs10[4]+0.04, yy10+0.03, col_ws10[4]-0.08, 0.22, size=8.5, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════
    # SLIDE 11 — ANALISIS FINAL
    # ══════════════════════════════════════════════════
    s11 = add_slide_header(prs, "Analisis Final — Lectura del Problema", periodo_label)
    bloques11 = [
        ("El proceso de control existe y funciona.",
         "Los desvios no tienden hacia cero sino que fluctuan de forma irregular, con picos recurrentes.", BLUE),
        ("Se corrigen los desvios?",
         "No de forma sostenida. Los mismos delegados y sectores reaparecen con desvios en distintos meses — las correcciones son temporales.", ORANGE),
        ("Riesgo de no actuar",
         "Sin seguimiento continuo, se consolida la normalizacion del desvio. Los delegados criticos perderan el objetivo principal: el trabajo en la linea.", RED),
    ]
    for i,(tit11,txt11,col11) in enumerate(bloques11):
        yy11 = 1.12 + i*1.4
        add_rect(s11, 0.4, yy11, 9.2, 1.25, LGRAY, col11, Pt(1))
        add_text(s11, tit11, 0.6, yy11+0.08, 8.8, 0.32, size=10.5, bold=True, color=col11)
        add_text(s11, txt11, 0.6, yy11+0.42, 8.8, 0.7, size=9.5, color=BLACK)

    # ══════════════════════════════════════════════════
    # SLIDE 12 — CONCLUSIONES Y PROXIMOS PASOS
    # ══════════════════════════════════════════════════
    s12 = prs.slides.add_slide(blank)
    add_rect(s12, 0, 0, 10, 5.625, NAVY)
    add_rect(s12, 0, 4.1, 10, 1.525, BLUE)
    add_text(s12, "CONCLUSIONES Y PROXIMOS PASOS", 0.5, 0.2, 9, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s12, 3, 0.78, 4, 0.04, ORANGE)

    add_text(s12, "CONCLUSIONES", 0.4, 0.95, 5.5, 0.28, size=9.5, bold=True, color=ORANGE)
    pico = max(D["tendencia"], key=lambda x:x["pct"]) if D["tendencia"] else {"label":"","pct":0}
    top_turno = D["turnos"][0] if D["turnos"] else {"turno":"","pct":0}
    top3secs = ", ".join(s["sector"].replace(" Division","") for s in D["sectores"][:3])
    top3pct  = fmt_pct(D["top5"][2]["pct"]) if len(D["top5"])>=3 else "N/A"
    concs12 = [
        f"El proceso de control genero identificar donde estan los desvios.",
        f"La tasa global es del {fmt_pct(D['tasa'])} — los desvios no se corrigen de forma sostenida.",
        f"Pico critico: {pico['label']} ({fmt_pct(pico['pct'])}%).",
        f"El turno {top_turno['turno']} concentra el mayor porcentaje ({fmt_pct(top_turno['pct'])}).",
        f"Sectores {top3secs} concentran el riesgo mas alto.",
        f"Los 3 delegados del Top 3 superan el {top3pct} individual sin correccion visible.",
    ]
    for i,c in enumerate(concs12):
        add_text(s12, f"✓  {c}", 0.4, 1.28+i*0.5, 5.5, 0.45, size=8.5, color="CADCFC")

    add_text(s12, "PROXIMOS PASOS", 6.3, 0.95, 3.4, 0.28, size=9.5, bold=True, color=ORANGE)
    pasos12 = [
        "Reunion con CIR y sectores criticos.",
        "Seguimiento individual de delegados atipicos.",
        "Reportar casos de LLT reincidentes.",
        "Monitoreo semanal de tasa de desvio.",
        "Generar nota en casos de reincidencia.",
    ]
    for i,p in enumerate(pasos12):
        add_rect(s12, 6.3, 1.32+i*0.62, 0.32, 0.32, ORANGE)
        add_text(s12, str(i+1), 6.3, 1.32+i*0.62, 0.32, 0.32, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s12, p, 6.7, 1.35+i*0.62, 3.0, 0.52, size=8.5, color="CADCFC")

    # ── Serializar a bytes ─────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
